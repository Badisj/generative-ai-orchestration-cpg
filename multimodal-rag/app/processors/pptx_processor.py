# FILE: app/processors/pptx_processor.py
"""
PowerPoint processor using python-pptx.
Extracts slide text with structure.
"""

from pathlib import Path
from typing import List
import logging

from app.processors.base import BaseProcessor, Document

logger = logging.getLogger(__name__)


class PowerPointProcessor(BaseProcessor):
    """
    Process PowerPoint presentations.
    
    Features:
    - Extracts text from all slides
    - Preserves slide titles and bullet structure
    - Handles speaker notes
    - Extracts table content
    """
    
    supported_extensions = [".pptx", ".ppt"]
    
    def __init__(self, include_notes: bool = True):
        """
        Initialize PowerPoint processor.
        
        Parameters
        ----------
        include_notes : bool
            Whether to include speaker notes.
        """
        self.include_notes = include_notes
    
    def can_process(self, file_path: Path) -> bool:
        """Check if file is a PowerPoint."""
        return file_path.suffix.lower() in self.supported_extensions
    
    def process(self, file_path: Path) -> List[Document]:
        """
        Extract text from PowerPoint presentation.
        
        Parameters
        ----------
        file_path : Path
            Path to PowerPoint file.
            
        Returns
        -------
        List[Document]
            List of documents, one per slide.
        """
        from pptx import Presentation
        from pptx.util import Inches
        
        file_path = Path(file_path)
        documents = []
        
        # Handle .ppt (older format)
        if file_path.suffix.lower() == ".ppt":
            logger.warning(f"Old .ppt format not fully supported: {file_path}")
            return []
        
        try:
            prs = Presentation(str(file_path))
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                content_parts = []
                
                # Extract slide title
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()
                    if title:
                        content_parts.append(f"# Slide {slide_num}: {title}")
                else:
                    content_parts.append(f"# Slide {slide_num}")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    # Skip title (already extracted)
                    if shape == slide.shapes.title:
                        continue
                    
                    # Text frames
                    if hasattr(shape, "text_frame"):
                        text = self._extract_text_frame(shape.text_frame)
                        if text:
                            content_parts.append(text)
                    
                    # Tables
                    if shape.has_table:
                        table_md = self._extract_table(shape.table)
                        if table_md:
                            content_parts.append(table_md)
                
                # Extract speaker notes
                if self.include_notes and slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        content_parts.append(f"\n**Speaker Notes:**\n{notes}")
                
                content = "\n\n".join(content_parts)
                
                if content.strip():
                    doc = self._create_document(
                        content=content,
                        source_file=file_path.name,
                        file_type="pptx",
                        page_number=slide_num,  # Using page_number for slide number
                        slide_title=slide.shapes.title.text if slide.shapes.title else None,
                    )
                    documents.append(doc)
                    
        except Exception as e:
            logger.error(f"Error processing PowerPoint {file_path}: {e}")
        
        logger.info(f"Extracted {len(documents)} slides from {file_path.name}")
        return documents
    
    def _extract_text_frame(self, text_frame) -> str:
        """
        Extract text from a text frame with bullet preservation.
        
        Parameters
        ----------
        text_frame : pptx.text.text.TextFrame
            PowerPoint text frame.
            
        Returns
        -------
        str
            Extracted text with bullet formatting.
        """
        lines = []
        
        for para in text_frame.paragraphs:
            text = para.text.strip()
            if text:
                # Add bullet point based on level
                level = para.level
                if level > 0:
                    indent = "  " * level
                    text = f"{indent}- {text}"
                elif para.level == 0 and len(text) < 100:
                    # Short top-level text might be a subtitle
                    text = f"- {text}"
                lines.append(text)
        
        return "\n".join(lines)
    
    def _extract_table(self, table) -> str:
        """
        Extract table as markdown.
        
        Parameters
        ----------
        table : pptx.table.Table
            PowerPoint table.
            
        Returns
        -------
        str
            Markdown formatted table.
        """
        rows = []
        
        for i, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                text = cell.text.strip().replace("\n", " ").replace("|", "\\|")
                cells.append(text)
            rows.append("| " + " | ".join(cells) + " |")
            
            # Header separator after first row
            if i == 0:
                rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
        
        return "\n".join(rows) if rows else ""
